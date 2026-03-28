from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor

OUT_PATH = r"d:\plant disease dataset\Plant_Disease_System_Architecture.pptx"


def add_box(slide, x, y, w, h, text, fill_rgb, font_size=12, bold=False, line_rgb=(80, 80, 80)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(20, 20, 20)
    return shape


def add_label(slide, x, y, w, h, text, size=11, bold=False, color=(40, 40, 40)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return box


def add_arrow(slide, x1, y1, x2, y2, color=(70, 70, 70), width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(*color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = add_label(slide, 0.35, 0.15, 12.8, 0.6, "System Architecture - Plant Disease Detection", size=34, bold=True, color=(40, 84, 52))

    # Outer architecture boundary
    boundary = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.2))
    boundary.fill.background()
    boundary.line.color.rgb = RGBColor(90, 90, 90)
    boundary.line.width = Pt(1.4)
    boundary.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # Left-side user and client
    user = add_box(slide, 0.8, 2.15, 1.25, 0.75, "Farmer/User", (230, 239, 248), font_size=12, bold=True)
    client = add_box(slide, 2.35, 2.15, 1.95, 0.75, "Web / Mobile\nInterface", (221, 236, 245), font_size=12, bold=True)

    # API and preprocessing
    api = add_box(slide, 4.75, 1.55, 2.0, 0.75, "Inference API\n(FastAPI/Flask)", (216, 230, 242), font_size=12, bold=True)

    pre_header = add_box(slide, 4.55, 2.55, 2.4, 0.55, "Preprocessing Pipeline", (212, 242, 229), font_size=12, bold=True)
    resize = add_box(slide, 4.65, 3.25, 1.05, 0.5, "Resize\n224x224", (232, 248, 238), font_size=10)
    to_tensor = add_box(slide, 5.85, 3.25, 1.05, 0.5, "ToTensor", (232, 248, 238), font_size=10)
    norm = add_box(slide, 5.25, 3.95, 1.05, 0.5, "Normalize\n(ImageNet)", (232, 248, 238), font_size=10)

    # Inference block
    model = add_box(slide, 7.35, 1.55, 2.45, 1.2, "Model Inference\nResNet50 (Transfer Learning)\nSoftmax + Top-k", (238, 225, 242), font_size=12, bold=True)

    # Explainability and knowledge modules
    gradcam = add_box(slide, 7.35, 3.1, 2.45, 0.95, "Explainability Module\nGrad-CAM Heatmap", (224, 235, 252), font_size=12, bold=True)
    tips = add_box(slide, 10.2, 3.1, 2.1, 0.95, "Recommendation Engine\nprevention_tips.py", (255, 235, 214), font_size=12, bold=True)

    # Output block
    response = add_box(slide, 10.2, 1.45, 2.1, 1.4, "Response Packaging\n- Predicted class\n- Confidence\n- Top-3 classes\n- Tips text\n- Heatmap", (255, 246, 205), font_size=11, bold=True)

    ui_output = add_box(slide, 10.2, 4.4, 2.1, 1.2, "Output to User\nDiagnosis + Advice\nVisual Explanation", (224, 245, 232), font_size=12, bold=True)

    # Optional model comparison block
    compare = add_box(slide, 7.35, 4.45, 2.45, 1.15, "Optional Evaluation\ncompare_models.py\nCustom CNN vs ResNet50", (244, 239, 255), font_size=11)

    # Section labels
    add_label(slide, 4.65, 6.35, 3.8, 0.3, "Online Inference Path", size=10, bold=True, color=(55, 55, 55))
    add_label(slide, 8.0, 6.35, 4.1, 0.3, "Explainability + Agronomic Guidance", size=10, bold=True, color=(55, 55, 55))

    # Arrows for primary flow
    add_arrow(slide, 2.05, 2.53, 2.35, 2.53)   # user -> client
    add_arrow(slide, 4.30, 2.53, 4.75, 1.92)   # client -> api
    add_arrow(slide, 5.75, 2.30, 5.75, 2.55)   # api -> preprocessing
    add_arrow(slide, 5.7, 3.5, 5.85, 3.5)      # resize -> tensor
    add_arrow(slide, 6.35, 3.75, 5.85, 3.95)   # tensor -> normalize
    add_arrow(slide, 6.3, 4.2, 7.35, 2.2)      # preprocess -> model
    add_arrow(slide, 9.80, 2.2, 10.2, 2.2)     # model -> response
    add_arrow(slide, 11.25, 2.85, 11.25, 4.4)  # response -> output
    add_arrow(slide, 10.2, 5.0, 2.05, 2.9)     # output -> user feedback

    # Arrows for explainability and tips
    add_arrow(slide, 8.55, 2.75, 8.55, 3.1)    # model -> gradcam
    add_arrow(slide, 9.80, 3.58, 10.2, 3.58)   # gradcam -> tips
    add_arrow(slide, 11.25, 3.1, 11.25, 2.85)  # tips -> response

    # Optional comparison relation
    add_arrow(slide, 8.55, 2.75, 8.55, 4.45)   # model -> compare

    # Footer
    add_label(slide, 0.55, 7.1, 12.2, 0.25, "AgriXAI architecture: upload -> preprocess -> ResNet50 inference -> Grad-CAM -> prevention tips -> response", size=9, color=(90, 90, 90))

    prs.save(OUT_PATH)
    print(f"Created: {OUT_PATH}")


if __name__ == "__main__":
    main()
