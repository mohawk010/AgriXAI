"""
generate_research_ppt.py
Generates the AgriXAI Plant Disease Classification research PPT
following the same 20-slide structure as Research.pptx (Major Project Review-II).

Usage:
    python scripts/generate_research_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
import copy
from lxml import etree


# ─────────────────────────────────────────────────────────────────
# THEME COLORS
# ─────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
C_ACCENT = RGBColor(0x16, 0x21, 0x3E)   # mid-blue
C_GREEN  = RGBColor(0x2E, 0x8B, 0x57)   # forest green (agri)
C_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_MUTED  = RGBColor(0xCC, 0xCC, 0xCC)   # light grey
C_GOLD   = RGBColor(0xF0, 0xA5, 0x00)   # amber accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=C_LIGHT,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=17, title=None, title_size=20):
    """Add a text box with optional title + bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = C_GOLD

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = C_LIGHT

    return txBox


def add_placeholder_image(slide, label, left, top, width, height):
    """Add a grey rectangle as a placeholder image with a label."""
    rect = add_rect(slide, left, top, width, height, RGBColor(0x3A, 0x3A, 0x5A))
    add_text_box(slide, f"[ {label} ]",
                 left, top + height // 3, width, Inches(0.5),
                 font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    return rect


def footer_bar(slide, slide_num, event_text="Major Project Review-II - 16th March 2026"):
    """Add bottom footer bar."""
    add_rect(slide, 0, SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.45), C_ACCENT)
    add_text_box(slide, event_text,
                 Inches(0.2), SLIDE_H - Inches(0.42), Inches(10), Inches(0.4),
                 font_size=11, color=C_MUTED, italic=True)
    add_text_box(slide, str(slide_num),
                 SLIDE_W - Inches(0.6), SLIDE_H - Inches(0.42), Inches(0.4), Inches(0.4),
                 font_size=11, color=C_GOLD, align=PP_ALIGN.RIGHT)


def slide_header(slide, title_text):
    """Add a green accent bar + slide title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_ACCENT)
    add_rect(slide, 0, Inches(1.1), Inches(0.12), Inches(5.9), C_GREEN)
    add_text_box(slide, title_text,
                 Inches(0.25), Inches(0.15),
                 Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=C_LIGHT)


# ─────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide."""
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)

    # Top accent strip
    add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), C_GREEN)

    # Logo placeholder (top-right)
    add_placeholder_image(sl, "SRM Logo", SLIDE_W - Inches(2.2), Inches(0.15), Inches(2.0), Inches(1.1))

    # Uni name
    add_text_box(sl, "SRM INSTITUTE OF SCIENCE AND TECHNOLOGY",
                 Inches(0.4), Inches(0.2), Inches(10), Inches(0.5),
                 font_size=15, bold=True, color=C_GOLD, align=PP_ALIGN.LEFT)
    add_text_box(sl, "DEPARTMENT OF COMPUTER SCIENCE AND ENGINEERING",
                 Inches(0.4), Inches(0.6), Inches(10), Inches(0.4),
                 font_size=12, color=C_MUTED, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.04), C_GREEN)

    # Main title
    add_text_box(sl,
                 "AgriXAI: Plant Disease Classification\nUsing Deep Learning & Explainable AI",
                 Inches(0.4), Inches(1.4), Inches(12.5), Inches(1.8),
                 font_size=36, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Leaf placeholder image
    add_placeholder_image(sl, "Plant / Leaf Image", Inches(5.5), Inches(3.3), Inches(2.4), Inches(1.6))

    # Metadata block
    meta = [
        ("Guide Name:",      "Dr. / Mr. [Guide Name]"),
        ("Designation:",     "Assistant Professor"),
        ("Department:",      "Computer Science & Engineering"),
        ("Batch ID:",        "[Batch ID]"),
        ("Reg. No:",         "[Registration Number]"),
        ("Name:",            "[Student Name]"),
    ]
    y = Inches(3.2)
    for label, val in meta:
        add_text_box(sl, label, Inches(0.5), y, Inches(1.8), Inches(0.35),
                     font_size=12, bold=True, color=C_GOLD)
        add_text_box(sl, val,   Inches(2.3), y, Inches(5.0), Inches(0.35),
                     font_size=12, color=C_MUTED)
        y += Inches(0.35)

    add_text_box(sl, "Major Project Review-II  |  16th March 2026",
                 Inches(0.4), SLIDE_H - Inches(0.9), Inches(12.5), Inches(0.4),
                 font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)
    add_rect(sl, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), C_GREEN)


def slide_02_overview(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Project Overview")
    footer_bar(sl, 2)

    rows = [
        ("Background:",  "Rapid rise in global food insecurity due to undetected plant diseases."),
        ("Problem:",     "Farmers lack affordable, fast, scalable disease diagnosis tools."),
        ("Motivation:",  "Need for an accurate, explainable AI system to aid precision agriculture."),
        ("Domain:",      "Agricultural AI, Computer Vision & Explainable AI (XAI)."),
    ]
    y = Inches(1.4)
    for label, val in rows:
        add_text_box(sl, label, Inches(0.55), y, Inches(2.0), Inches(0.45),
                     font_size=16, bold=True, color=C_GOLD)
        add_text_box(sl, val,   Inches(2.55), y, Inches(10.2), Inches(0.45),
                     font_size=16, color=C_LIGHT)
        add_rect(sl, Inches(0.55), y + Inches(0.42), Inches(12.3), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.72)


def slide_03_problem(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Problem Statement")
    footer_bar(sl, 3)

    bullets = [
        "Farmers visually diagnose diseases — slow, error-prone, expertise-dependent.",
        "Most deep learning systems are opaque 'black boxes' with no interpretability.",
        "High accuracy alone insufficient for trust in high-stakes agricultural decisions.",
        "Lack of accessible tools combining classification + visual reasoning for field use.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.4), Inches(8.5), Inches(4.5),
                   font_size=18)
    add_placeholder_image(sl, "Problem Illustration\n(e.g. diseased leaf vs healthy)",
                          Inches(9.2), Inches(1.4), Inches(3.8), Inches(4.5))


def slide_04_objectives(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Project Objectives")
    footer_bar(sl, 4)

    add_text_box(sl, "Goal: Build an accurate, interpretable plant disease detection system.",
                 Inches(0.55), Inches(1.35), Inches(12.5), Inches(0.45),
                 font_size=17, bold=True, color=C_GOLD)
    bullets = [
        "Train a Custom CNN (5-block) from scratch on 87k+ plant leaf images.",
        "Fine-tune ResNet50 (ImageNet pretrained) via transfer learning.",
        "Implement Grad-CAM heatmaps for visual explainability of predictions.",
        "Provide agronomic prevention & treatment tips per disease class.",
        "Compare both models on accuracy, loss, and inference efficiency.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.9), Inches(12.3), Inches(4.5),
                   font_size=17)


def slide_05_literature(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Literature Review")
    footer_bar(sl, 5)

    bullets = [
        "Reviewed 10+ papers on CNN-based plant disease classification.",
        "Current SOTA: EfficientNet, Vision Transformers — high accuracy but low interpretability.",
        "PlantVillage dataset widely benchmarked; augmented versions push accuracy limits.",
        "Gap: Few systems combine Grad-CAM XAI with practical agronomic advisory outputs.",
        "Transfer learning (ResNet, VGG, Inception) consistently outperforms CNNs from scratch.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.35), Inches(12.3), Inches(4.8),
                   font_size=17)


def slide_06_architecture(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "System Architecture")
    footer_bar(sl, 6)

    add_placeholder_image(sl, "System Architecture Diagram\n(Input → Preprocessing → CNN/ResNet → Grad-CAM → Tips Output)",
                          Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.2))


def slide_07_proposed(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Proposed System")
    footer_bar(sl, 7)

    add_text_box(sl, "Proposed Solution: Dual-Model Explainable Plant Disease Classifier",
                 Inches(0.55), Inches(1.35), Inches(12.3), Inches(0.45),
                 font_size=17, bold=True, color=C_GOLD)
    bullets = [
        "Architecture 1: Custom CNN — 5 convolutional blocks (32→512 filters), Dropout, Dense(38).",
        "Architecture 2: ResNet50 — ImageNet pretrained; conv5 block fine-tuned, GlobalAvgPool, Dropout(0.5).",
        "XAI Layer: Grad-CAM overlays highlight discriminative leaf regions per prediction.",
        "Advisory Layer: Per-class prevention & treatment tips provided post-inference via prevention_tips.py.",
        "Inference: Single-image pipeline (predict.py) returns class, confidence, top-3, heatmap, and tips.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.9), Inches(12.3), Inches(4.5),
                   font_size=17)


def slide_08_dataset(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Dataset Description")
    footer_bar(sl, 8)

    rows = [
        ("Source:",   "New Plant Diseases Dataset (Augmented) — Kaggle."),
        ("Size:",     "~87,000 training images | ~22,000 validation images."),
        ("Classes:",  "38 disease/healthy categories across 14 plant species."),
        ("Format:",   "128×128 RGB (CNN) | 224×224 RGB (ResNet50); ImageFolder structure."),
        ("Split:",    "Train / Validation folders; augmentation includes flip, rotation."),
    ]
    y = Inches(1.4)
    for label, val in rows:
        add_text_box(sl, label, Inches(0.55), y, Inches(2.0), Inches(0.45),
                     font_size=16, bold=True, color=C_GOLD)
        add_text_box(sl, val,   Inches(2.55), y, Inches(6.8), Inches(0.45),
                     font_size=16, color=C_LIGHT)
        add_rect(sl, Inches(0.55), y + Inches(0.42), Inches(9.5), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.72)
    add_placeholder_image(sl, "Dataset Sample Images\n(Healthy vs. Diseased Leaves)", Inches(9.6), Inches(1.4), Inches(3.5), Inches(4.5))


def slide_09_methodology(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Methodology")
    footer_bar(sl, 9)

    steps = [
        ("1. Input",        "Leaf image uploaded; resized to 128×128 or 224×224."),
        ("2. Augmentation", "RandomHorizontalFlip, RandomRotation(15°) applied during training."),
        ("3. Model Forward","Image passed through CNN / ResNet50; logits computed."),
        ("4. Prediction",   "Softmax → top-1 class + confidence; top-3 alternatives surfaced."),
        ("5. Grad-CAM",     "Gradient-weighted activation map overlaid on original image."),
        ("6. Advisory",     "prevention_tips.py returns formatted treatment & prevention text."),
    ]
    y = Inches(1.35)
    for step, desc in steps:
        add_text_box(sl, step, Inches(0.55), y, Inches(2.4), Inches(0.4),
                     font_size=15, bold=True, color=C_GOLD)
        add_text_box(sl, desc, Inches(3.0), y, Inches(9.8), Inches(0.4),
                     font_size=15, color=C_LIGHT)
        add_rect(sl, Inches(0.55), y + Inches(0.38), Inches(12.3), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.62)


def slide_10_tools(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Tools & Technologies")
    footer_bar(sl, 10)

    cols = [
        ("Deep Learning (TF/Keras)",
         ["TensorFlow 2.x", "Keras (Custom CNN + ResNet50)", "EarlyStopping & ModelCheckpoint callbacks"]),
        ("Transfer Learning (PyTorch)",
         ["PyTorch 2.0+", "torchvision ResNet50 (ImageNet weights)", "Adam optimizer, CrossEntropyLoss"]),
        ("XAI & Visualization",
         ["Grad-CAM (custom gradcam.py)", "Matplotlib / Seaborn", "OpenCV (heatmap overlay)"]),
        ("Data & Utilities",
         ["scikit-learn (metrics, confusion matrix)", "NumPy / PIL", "Jupyter Notebook"]),
    ]
    xs = [Inches(0.4), Inches(3.5), Inches(6.65), Inches(9.8)]
    for (title, items), x in zip(cols, xs):
        add_rect(sl, x, Inches(1.25), Inches(2.9), Inches(4.9), C_ACCENT)
        add_text_box(sl, title, x + Inches(0.1), Inches(1.3), Inches(2.7), Inches(0.45),
                     font_size=13, bold=True, color=C_GREEN)
        y2 = Inches(1.8)
        for item in items:
            add_text_box(sl, f"• {item}", x + Inches(0.15), y2, Inches(2.6), Inches(0.4),
                         font_size=13, color=C_LIGHT)
            y2 += Inches(0.42)


def slide_11_modules(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Modules Implemented")
    footer_bar(sl, 11)

    modules = [
        ("Dataset Module (utils/dataset.py)",
         ["ImageFolder-based dataloaders", "Train/val augmentation pipeline", "Configurable img_size, batch_size, workers"]),
        ("Model Module (models/resnet50.py)",
         ["ResNet50 with selective layer freezing", "Dropout + custom FC head (38 classes)", "Trainable_layers param for fine-tune control"]),
        ("Training Module (training/train_resnet.py)",
         ["Epoch loop with train/val phases", "Best-model checkpoint saving (val_acc)", "Loss + accuracy curve plotting"]),
        ("Prediction & XAI (utils/predict.py)",
         ["Single-image inference pipeline", "Top-K class prediction with confidence", "Grad-CAM heatmap generation (optional)"]),
        ("Advisory Module (utils/prevention_tips.py)",
         ["38-class agronomic tips dictionary", "Formatted text output per disease", "Integrated into inference pipeline"]),
    ]
    y = Inches(1.3)
    for title, bullets in modules:
        add_text_box(sl, title, Inches(0.55), y, Inches(12.3), Inches(0.35),
                     font_size=14, bold=True, color=C_GOLD)
        line = "  |  ".join(bullets)
        add_text_box(sl, f"   {line}", Inches(0.55), y + Inches(0.33), Inches(12.3), Inches(0.35),
                     font_size=13, color=C_MUTED)
        add_rect(sl, Inches(0.55), y + Inches(0.68), Inches(12.3), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.85)


def slide_12_screens(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Implementation Screens")
    footer_bar(sl, 12)

    imgs = [
        ("Training Loss & Accuracy Curves", Inches(0.4), Inches(1.3), Inches(4.0), Inches(3.0)),
        ("Confusion Matrix Heatmap",        Inches(4.6), Inches(1.3), Inches(4.0), Inches(3.0)),
        ("Grad-CAM Overlay Output",         Inches(9.2), Inches(1.3), Inches(3.9), Inches(3.0)),
        ("Sample Prediction Report",        Inches(0.4), Inches(4.5), Inches(5.5), Inches(2.0)),
        ("Disease Class Distribution",      Inches(6.1), Inches(4.5), Inches(6.8), Inches(2.0)),
    ]
    for label, l, t, w, h in imgs:
        add_placeholder_image(sl, label, l, t, w, h)


def slide_13_results(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Experimental Results")
    footer_bar(sl, 13)

    # Table-style results
    headers = ["Metric", "Custom CNN", "ResNet50"]
    rows_data = [
        ["Final Train Accuracy", "98.12%", "99.81%"],
        ["Final Val Accuracy",   "96.26%", "99.45%"],
        ["Final Train Loss",     "0.0571",  "0.0059"],
        ["Final Val Loss",       "0.1227",  "0.0214"],
    ]
    col_x = [Inches(0.55), Inches(4.8), Inches(8.8)]
    col_w = [Inches(4.0),  Inches(3.8), Inches(3.8)]
    # Header
    y = Inches(1.35)
    for (h, x, w) in zip(headers, col_x, col_w):
        add_rect(sl, x, y, w, Inches(0.42), C_GREEN)
        add_text_box(sl, h, x + Inches(0.1), y + Inches(0.05), w - Inches(0.15), Inches(0.35),
                     font_size=15, bold=True, color=C_LIGHT)
    y += Inches(0.42)
    for i, row in enumerate(rows_data):
        row_color = C_ACCENT if i % 2 == 0 else RGBColor(0x22, 0x22, 0x3D)
        for val, x, w in zip(row, col_x, col_w):
            add_rect(sl, x, y, w, Inches(0.42), row_color)
            add_text_box(sl, val, x + Inches(0.1), y + Inches(0.05), w - Inches(0.1), Inches(0.35),
                         font_size=14, color=C_LIGHT)
        y += Inches(0.42)

    add_text_box(sl, "ResNet50 (transfer learning) significantly outperforms Custom CNN.",
                 Inches(0.55), y + Inches(0.2), Inches(12.3), Inches(0.4),
                 font_size=15, bold=True, color=C_GOLD)
    add_placeholder_image(sl, "Accuracy / Loss Curve Plot", Inches(0.55), y + Inches(0.7), Inches(12.0), Inches(2.5))


def slide_14_discussion(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Result Discussion")
    footer_bar(sl, 14)

    bullets = [
        "ResNet50 (val acc 99.45%) drastically outperforms Custom CNN (96.26%) due to ImageNet pretraining.",
        "Transfer learning with selective layer unfreezing (conv5 block) enables fast convergence.",
        "Grad-CAM heatmaps confirm model focuses on correct leaf lesion areas — improving trust.",
        "Custom CNN still competitive for lightweight edge deployment (lower model size).",
        "Overfitting mitigated via Dropout (0.25 / 0.4 CNN; 0.5 ResNet) and data augmentation.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.35), Inches(12.3), Inches(4.8), font_size=17)


def slide_15_work_done(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Work Completed So Far")
    footer_bar(sl, 15)

    phases = [
        ("Phase 1", "Dataset Setup & EDA", "Downloaded 87k+ images, organized ImageFolder structure, exploratory class-balance analysis."),
        ("Phase 2", "Custom CNN Training",  "Built 5-block CNN in TF/Keras; trained with Adam (lr=0.0001); achieved 96.26% val accuracy."),
        ("Phase 3", "ResNet50 Fine-Tuning", "ImageNet pretrained ResNet50; unfroze conv5 block; achieved 99.45% val accuracy."),
        ("Phase 4", "Grad-CAM Integration", "Implemented gradcam.py; integrated into predict.py for visual XAI on single-image inference."),
        ("Phase 5", "Advisory System",      "Built 38-class prevention_tips.py with agronomic treatment advice per disease category."),
    ]
    y = Inches(1.3)
    for tag, title, desc in phases:
        add_rect(sl, Inches(0.4), y, Inches(1.0), Inches(0.5), C_GREEN)
        add_text_box(sl, tag, Inches(0.4), y + Inches(0.05), Inches(1.0), Inches(0.4),
                     font_size=13, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
        add_text_box(sl, title, Inches(1.55), y, Inches(2.8), Inches(0.5),
                     font_size=14, bold=True, color=C_GOLD)
        add_text_box(sl, desc, Inches(4.5), y, Inches(8.6), Inches(0.5),
                     font_size=13, color=C_MUTED)
        y += Inches(0.72)


def slide_16_publication(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Publication Status")
    footer_bar(sl, 16)

    rows = [
        ("Status:",        "Drafting manuscript — methodology and experimental results sections in progress."),
        ("Target Venue:",  "IEEE / Springer journal focused on agricultural AI, precision farming, or XAI."),
        ("Working Title:", '"AgriXAI: Explainable Deep Learning for Multi-Class Plant Disease Detection using Grad-CAM"'),
        ("Scope:",         "Novel contribution: combining transfer learning + XAI + agronomic advisory in a single pipeline."),
    ]
    y = Inches(1.5)
    for label, val in rows:
        add_text_box(sl, label, Inches(0.55), y, Inches(2.2), Inches(0.5),
                     font_size=16, bold=True, color=C_GOLD)
        add_text_box(sl, val, Inches(2.8), y, Inches(10.0), Inches(0.5),
                     font_size=16, color=C_LIGHT)
        add_rect(sl, Inches(0.55), y + Inches(0.47), Inches(12.3), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.82)


def slide_17_challenges(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Challenges Faced")
    footer_bar(sl, 17)

    chals = [
        ("Dataset Size",       "Managing ~1.3 GB dataset; GPU memory constraints during ResNet50 training with batch_size=32."),
        ("Grad-CAM Alignment", "Hooking correct target layers across TF (Keras) and PyTorch ResNet architectures separately."),
        ("Class Imbalance",    "38-class label visualization; long-tail classes requiring careful confusion-matrix analysis."),
        ("Augmentation Tuning","Balancing RandomRotation / Flip without introducing harmful distribution shift."),
    ]
    y = Inches(1.35)
    for tag, desc in chals:
        add_rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.5), C_GREEN)
        add_text_box(sl, tag,  Inches(0.6), y, Inches(2.8), Inches(0.5),
                     font_size=15, bold=True, color=C_GOLD)
        add_text_box(sl, desc, Inches(3.5), y, Inches(9.5), Inches(0.5),
                     font_size=15, color=C_LIGHT)
        add_rect(sl, Inches(0.55), y + Inches(0.5), Inches(12.3), Inches(0.02), RGBColor(0x33,0x33,0x55))
        y += Inches(0.85)


def slide_18_conclusion(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "Conclusion")
    footer_bar(sl, 18)

    bullets = [
        "Achieved 99.45% validation accuracy with ResNet50 transfer learning on 38-class plant disease dataset.",
        "Demonstrated that Custom CNN (96.26%) provides a strong lightweight alternative.",
        "Grad-CAM integration eliminates the 'black box' problem — predictions are visually interpretable.",
        "Prevention tips advisory system bridges AI output to real agronomic action for farmers.",
        "Scalable pipeline with separate training, inference, and comparison modules.",
    ]
    add_bullet_box(sl, bullets, Inches(0.55), Inches(1.35), Inches(12.3), Inches(4.8), font_size=17)


def slide_19_references(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    slide_header(sl, "References")
    footer_bar(sl, 19)

    refs = [
        "Hughes, D. & Salathé, M. An open access repository of images on plant health to enable the development of mobile disease diagnostics. arXiv, 2015.",
        "He, K., et al. Deep Residual Learning for Image Recognition. CVPR, 2016.",
        "Selvaraju, R. R., et al. Grad-CAM: Visual Explanations from Deep Networks via Gradient-based Localization. ICCV, 2017.",
        "Mohanty, S. P., et al. Using Deep Learning for Image-Based Plant Disease Detection. Frontiers in Plant Science, 2016.",
        "Goodfellow, I., Bengio, Y., & Courville, A. Deep Learning. MIT Press, 2016.",
    ]
    y = Inches(1.35)
    for i, ref in enumerate(refs, 1):
        add_text_box(sl, f"[{i}]  {ref}",
                     Inches(0.55), y, Inches(12.3), Inches(0.55),
                     font_size=14, color=C_LIGHT)
        y += Inches(0.58)


def slide_20_thankyou(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_DARK)
    add_rect(sl, 0, 0, SLIDE_W, Inches(0.12), C_GREEN)
    add_rect(sl, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), C_GREEN)

    add_text_box(sl, "Thank You",
                 Inches(0), Inches(2.3), SLIDE_W, Inches(1.4),
                 font_size=54, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(sl, "AgriXAI — Plant Disease Classification Using Deep Learning & Explainable AI",
                 Inches(0), Inches(3.7), SLIDE_W, Inches(0.5),
                 font_size=16, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

    add_placeholder_image(sl, "Healthy vs Diseased Leaf — Side by Side",
                          Inches(5.4), Inches(4.4), Inches(2.5), Inches(1.4))

    add_text_box(sl, "Major Project Review-II  |  16th March 2026",
                 Inches(0), SLIDE_H - Inches(0.7), SLIDE_W, Inches(0.4),
                 font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_problem(prs)
    slide_04_objectives(prs)
    slide_05_literature(prs)
    slide_06_architecture(prs)
    slide_07_proposed(prs)
    slide_08_dataset(prs)
    slide_09_methodology(prs)
    slide_10_tools(prs)
    slide_11_modules(prs)
    slide_12_screens(prs)
    slide_13_results(prs)
    slide_14_discussion(prs)
    slide_15_work_done(prs)
    slide_16_publication(prs)
    slide_17_challenges(prs)
    slide_18_conclusion(prs)
    slide_19_references(prs)
    slide_20_thankyou(prs)

    out_path = "AgriXAI_Research_PPT.pptx"
    prs.save(out_path)
    print(f"✅  Saved: {os.path.abspath(out_path)}")


if __name__ == "__main__":
    main()
