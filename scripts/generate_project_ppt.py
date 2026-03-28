from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = r"d:\plant disease dataset\Plant_Disease_Detection_20_Slides.pptx"

SLIDES = [
    {
        "title": "Plant Disease Detection Using Deep Learning",
        "subtitle": [
            "ResNet50 + Grad-CAM + Prevention Recommendations",
            "Name: ____________________",
            "Department/Institute: ____________________",
            "Date: ____________________",
        ],
    },
    {
        "title": "Agenda",
        "bullets": [
            "Introduction and problem statement",
            "Objectives and motivation",
            "Dataset overview",
            "Model architecture and training",
            "Results and explainability",
            "Deployment concept and future scope",
        ],
    },
    {
        "title": "Introduction",
        "bullets": [
            "Plant diseases significantly reduce crop quality and yield",
            "Early diagnosis helps reduce losses and improve farm decisions",
            "Manual diagnosis is time-consuming and often inconsistent",
            "Computer vision enables fast and automated leaf disease detection",
        ],
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Farmers may struggle to identify disease type accurately",
            "Several diseases have visually similar symptoms",
            "Field support may not be available in real time",
            "Need an accurate, scalable, and low-latency diagnosis assistant",
        ],
    },
    {
        "title": "Project Objectives",
        "bullets": [
            "Build a multiclass plant disease image classifier",
            "Use transfer learning with ResNet50",
            "Provide top-k class predictions with confidence",
            "Return prevention/treatment guidance for predicted disease",
            "Improve trust with Grad-CAM visual explanations",
        ],
    },
    {
        "title": "Motivation and Background",
        "bullets": [
            "Traditional methods rely heavily on handcrafted features",
            "Deep CNNs learn discriminative features automatically",
            "Transfer learning reduces compute cost and training time",
            "Goal: practical AI support for precision agriculture",
        ],
    },
    {
        "title": "Dataset Overview",
        "bullets": [
            "Dataset: New Plant Diseases Dataset (Augmented)",
            "Covers multiple crops and disease categories",
            "Directory structure follows class-wise folders",
            "Train and validation splits are already provided",
        ],
    },
    {
        "title": "Class Distribution",
        "bullets": [
            "Total classes: 38 (healthy + diseased)",
            "Crops include Apple, Tomato, Potato, Corn, Grape, etc.",
            "Balanced representation is important for generalization",
            "Class names are mapped from folder labels",
        ],
    },
    {
        "title": "Data Preprocessing",
        "bullets": [
            "Resize all images to 224 x 224",
            "Convert images to tensors",
            "Normalize using ImageNet mean and standard deviation",
            "Ensure preprocessing parity between training and inference",
        ],
    },
    {
        "title": "Model Architecture",
        "bullets": [
            "Backbone: ResNet50 pretrained on ImageNet",
            "Final classifier head adapted to 38 classes",
            "Forward pass outputs logits for each class",
            "Softmax converts logits to probabilities",
        ],
    },
    {
        "title": "Why ResNet50",
        "bullets": [
            "Residual blocks help prevent vanishing gradients",
            "Strong benchmark performance in visual recognition tasks",
            "Good trade-off between accuracy and complexity",
            "Well-supported for transfer learning workflows",
        ],
    },
    {
        "title": "Training Setup",
        "bullets": [
            "Loss function: Cross-Entropy",
            "Optimizer: Adam/SGD (as configured in training script)",
            "Validation after each epoch",
            "Best model checkpoint saved for inference",
            "Metrics logged using history JSON files",
        ],
    },
    {
        "title": "Training Workflow",
        "bullets": [
            "Load data using train/validation dataloaders",
            "Run epoch-wise train and validation loops",
            "Track loss and accuracy curves",
            "Prevent overfitting with monitoring and checkpointing",
        ],
    },
    {
        "title": "Evaluation Metrics",
        "bullets": [
            "Primary metric: validation accuracy",
            "Support metrics: train/val loss trends",
            "Top-k probabilities for confidence-aware decisions",
            "Optional class-wise analysis using confusion matrix",
        ],
    },
    {
        "title": "Model Performance",
        "bullets": [
            "Show final training and validation accuracy",
            "Insert accuracy-vs-epochs chart",
            "Insert loss-vs-epochs chart",
            "Discuss signs of underfitting/overfitting",
        ],
    },
    {
        "title": "Inference Pipeline",
        "bullets": [
            "Input image is preprocessed and sent to the model",
            "Predict top class, confidence, and top-k labels",
            "Fetch disease-specific prevention/treatment tips",
            "Return all outputs in a structured result dictionary",
        ],
    },
    {
        "title": "Explainability with Grad-CAM",
        "bullets": [
            "Generates class activation heatmaps",
            "Highlights image regions influencing prediction",
            "Verifies whether model focuses on infected areas",
            "Improves transparency and user trust",
        ],
    },
    {
        "title": "Sample Prediction Output",
        "bullets": [
            "Predicted class: e.g., Tomato___Late_blight",
            "Confidence score: e.g., 97.3%",
            "Top-3 classes with probabilities",
            "Human-readable prevention and treatment guidance",
        ],
    },
    {
        "title": "Limitations and Challenges",
        "bullets": [
            "Real field images may include noise and complex backgrounds",
            "Lighting and camera quality can impact performance",
            "Some disease classes have subtle visual differences",
            "Model should be evaluated on more diverse real-world data",
        ],
    },
    {
        "title": "Conclusion and Future Work",
        "bullets": [
            "A deep learning model can effectively classify plant leaf diseases",
            "System combines prediction, confidence, explanation, and guidance",
            "Future: mobile app deployment and multilingual recommendations",
            "Future: add new crops, real-time camera inference, and field validation",
        ],
    },
]


def style_title(shape):
    tf = shape.text_frame
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(25, 64, 36)


def add_subtitle_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()
    for i, line in enumerate(lines):
        p = subtitle.paragraphs[0] if i == 0 else subtitle.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if i == 0 else 16)
        p.font.color.rgb = RGBColor(55, 55, 55)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(20, 20, 20)


def main():
    prs = Presentation()

    # Widescreen aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, slide_data in enumerate(SLIDES):
        if idx == 0:
            add_subtitle_slide(prs, slide_data["title"], slide_data["subtitle"])
        else:
            add_bullet_slide(prs, slide_data["title"], slide_data["bullets"])

    prs.save(OUTPUT_PATH)
    print(f"Created presentation: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
